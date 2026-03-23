#!/usr/bin/env python3
"""Parse common business attachment types into extracted text, document type, key fields, and a concise summary.

Usage:
    python scripts/parse_attachment.py /path/to/file
    python scripts/parse_attachment.py /path/to/file --max-chars 12000
"""
from __future__ import annotations

import argparse
import csv
import json
import os
import re
import sys
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Dict, List, Tuple

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

TEXT_EXTS = {'.txt', '.md', '.json', '.eml', '.log', '.rtf'}
SHEET_EXTS = {'.xlsx', '.xlsm', '.csv', '.tsv'}
INVOICE_TERMS = ['invoice', 'bill to', 'payment terms', 'remit', 'amount due', 'tax invoice']
CONTRACT_TERMS = ['agreement', 'effective date', 'term', 'termination', 'renewal', 'msa', 'sow', 'nda']
RESUME_TERMS = ['experience', 'skills', 'education', 'curriculum vitae', 'resume', 'employment']
REPORT_TERMS = ['executive summary', 'recommendation', 'findings', 'metrics', 'analysis', 'report']
CURRENCY_PATTERN = r'(?:USD|EUR|GBP|ZAR|KES|NGN|R|\\$)\\s?[0-9][0-9,]*(?:\\.\\d{2})?'

@dataclass
class Result:
    file_name: str
    file_type: str
    detected_document_type: str
    extracted_fields: Dict[str, str]
    summary: str
    warnings: List[str]
    extracted_text_preview: str


def normalize_space(text: str) -> str:
    return re.sub(r'\s+', ' ', text).strip()


def limit(text: str, max_chars: int) -> str:
    text = text.strip()
    return text[:max_chars] + ('...' if len(text) > max_chars else '')


def read_text(path: Path) -> Tuple[str, List[str]]:
    warnings: List[str] = []
    suffix = path.suffix.lower()

    if suffix == '.pdf':
        if PdfReader is None:
            raise RuntimeError('missing dependency: pypdf')
        pages = []
        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages, start=1):
            try:
                pages.append(page.extract_text() or '')
            except Exception as exc:
                warnings.append(f'page {i} could not be read: {exc}')
        text = '\n\n'.join(pages)
        if not normalize_space(text):
            warnings.append('pdf contained little or no extractable text; it may be scanned or image-based')
        return text, warnings

    if suffix == '.docx':
        if Document is None:
            raise RuntimeError('missing dependency: python-docx')
        doc = Document(str(path))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n'.join(paras), warnings

    if suffix == '.pptx':
        if Presentation is None:
            raise RuntimeError('missing dependency: python-pptx')
        prs = Presentation(str(path))
        chunks: List[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            chunks.append(f'[Slide {idx}]')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    chunks.append(shape.text)
        return '\n'.join(chunks), warnings

    if suffix in {'.xlsx', '.xlsm'}:
        if load_workbook is None:
            raise RuntimeError('missing dependency: openpyxl')
        wb = load_workbook(str(path), read_only=True, data_only=True)
        chunks: List[str] = []
        for ws in wb.worksheets:
            chunks.append(f'[Sheet: {ws.title}]')
            for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                vals = ['' if v is None else str(v) for v in row]
                if any(v.strip() for v in vals):
                    chunks.append(' | '.join(vals))
                if row_idx >= 50:
                    chunks.append('[truncated after 50 rows]')
                    break
        return '\n'.join(chunks), warnings

    if suffix in {'.csv', '.tsv'}:
        delim = ',' if suffix == '.csv' else '\t'
        chunks: List[str] = []
        with path.open('r', encoding='utf-8', errors='replace', newline='') as fh:
            reader = csv.reader(fh, delimiter=delim)
            for idx, row in enumerate(reader, start=1):
                if any(cell.strip() for cell in row):
                    chunks.append(' | '.join(row))
                if idx >= 50:
                    chunks.append('[truncated after 50 rows]')
                    break
        return '\n'.join(chunks), warnings

    if suffix in TEXT_EXTS:
        return path.read_text(encoding='utf-8', errors='replace'), warnings

    raise ValueError(f'unsupported file type: {suffix or "no extension"}')


def match(pattern: str, text: str, flags: int = re.I) -> str:
    m = re.search(pattern, text, flags)
    return normalize_space(m.group(1)) if m else 'not stated'


def money_values(text: str) -> List[str]:
    vals = re.findall(CURRENCY_PATTERN, text, flags=re.I)
    return list(dict.fromkeys(vals))[:5]


def dates(text: str) -> List[str]:
    pats = [
        r'\b\d{4}-\d{2}-\d{2}\b',
        r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b',
        r'\b(?:jan|feb|mar|apr|may|jun|jul|aug|sep|sept|oct|nov|dec)[a-z]*\s+\d{1,2},\s*\d{4}\b',
    ]
    out: List[str] = []
    for pat in pats:
        out.extend(re.findall(pat, text, flags=re.I))
    return list(dict.fromkeys(out))[:8]


def classify(text: str, file_name: str) -> str:
    low = f'{file_name.lower()}\n{text[:10000].lower()}'
    scores = {
        'invoice_or_quote': sum(t in low for t in INVOICE_TERMS),
        'contract_or_legal_document': sum(t in low for t in CONTRACT_TERMS),
        'resume_or_candidate_profile': sum(t in low for t in RESUME_TERMS),
        'report_or_presentation': sum(t in low for t in REPORT_TERMS),
        'spreadsheet_or_data_export': 2 if Path(file_name).suffix.lower() in SHEET_EXTS else 0,
    }
    label, score = max(scores.items(), key=lambda kv: kv[1])
    return label if score > 0 else 'general_business_document'


def extract_invoice_fields(text: str) -> Dict[str, str]:
    return {
        'issuer': match(r'(?:from|supplier|vendor)[:\s]+([^\n]+)', text),
        'recipient': match(r'(?:bill to|invoice to)[:\s]+([^\n]+)', text),
        'invoice_number': match(r'(?:invoice\s*(?:number|no\.?|#))[:\s]+([A-Z0-9\-/]+)', text),
        'issue_date': match(r'(?:invoice date|issue date|date)[:\s]+([^\n]+)', text),
        'due_date': match(r'(?:due date|payment due)[:\s]+([^\n]+)', text),
        'currency_or_amounts': ', '.join(money_values(text)) or 'not stated',
        'payment_terms': match(r'(?:payment terms|terms)[:\s]+([^\n]+)', text),
        'purchase_order_reference': match(r'(?:po number|purchase order|po#)[:\s]+([^\n]+)', text),
    }


def extract_contract_fields(text: str) -> Dict[str, str]:
    return {
        'parties': match(r'between\s+(.+?)\s+and\s+(.+?)(?:\.|\n)', text),
        'effective_date': match(r'(?:effective date|commencement date)[:\s]+([^\n]+)', text),
        'term': match(r'(?:term)[:\s]+([^\n]+)', text),
        'renewal': match(r'(?:renewal|auto-renewal|extension)[:\s]+([^\n]+)', text),
        'termination': match(r'(?:termination|terminate)[:\s]+([^\n]+)', text),
        'fees_or_pricing': ', '.join(money_values(text)) or 'not stated',
        'signature_status': 'signed' if re.search(r'signature|signed by|executed by', text, re.I) else 'not stated',
    }


def extract_resume_fields(text: str) -> Dict[str, str]:
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return {
        'candidate_name': lines[0] if lines else 'not stated',
        'current_or_recent_employer': match(r'(?:experience|employment).*?\b(?:at|with)\s+([^\n,]+)', text, flags=re.I | re.S),
        'years_of_experience': match(r'(\d+\+?\s+years? of experience)', text),
        'location': match(r'(?:location|based in)[:\s]+([^\n]+)', text),
        'availability': match(r'(?:notice period|availability|available from)[:\s]+([^\n]+)', text),
        'key_skills': ', '.join(re.findall(r'\b(?:python|sql|excel|power bi|java|salesforce|aws|azure|project management|financial modeling)\b', text, re.I)[:8]) or 'not stated',
    }


def extract_report_fields(text: str) -> Dict[str, str]:
    return {
        'topic': match(r'(?:title|subject|topic)[:\s]+([^\n]+)', text),
        'headline_metrics': ', '.join(money_values(text) + dates(text)[:3]) or 'not stated',
        'recommendations_present': 'yes' if re.search(r'recommendation|next steps|action items', text, re.I) else 'not stated',
        'decisions_needed': 'yes' if re.search(r'decision required|approval|approve|sign[- ]?off', text, re.I) else 'not stated',
    }


def extract_generic_fields(text: str) -> Dict[str, str]:
    return {
        'dates_found': ', '.join(dates(text)) or 'not stated',
        'amounts_found': ', '.join(money_values(text)) or 'not stated',
        'reference_ids': ', '.join(re.findall(r'\b[A-Z]{2,6}-\d{2,}\b', text)[:10]) or 'not stated',
    }


def summarize(doc_type: str, fields: Dict[str, str], text: str) -> str:
    first_sentence = normalize_space(re.split(r'(?<=[.!?])\s+', text.strip())[0]) if text.strip() else ''
    if doc_type == 'invoice_or_quote':
        return f"This appears to be an invoice or quote. The issuer is {fields['issuer']}, the invoice number is {fields['invoice_number']}, and the due date is {fields['due_date']}. Amounts identified: {fields['currency_or_amounts']}."
    if doc_type == 'contract_or_legal_document':
        return f"This appears to be a contract or legal document. Effective date: {fields['effective_date']}. Term: {fields['term']}. Renewal: {fields['renewal']}. Fees or pricing found: {fields['fees_or_pricing']}."
    if doc_type == 'resume_or_candidate_profile':
        return f"This appears to be a candidate profile or resume. Candidate: {fields['candidate_name']}. Current or recent employer: {fields['current_or_recent_employer']}. Key skills: {fields['key_skills']}."
    if doc_type == 'report_or_presentation':
        return f"This appears to be a report or presentation. Topic: {fields['topic']}. Headline metrics or dates found: {fields['headline_metrics']}. Decisions needed: {fields['decisions_needed']}."
    return f"This appears to be a general business document. Opening content: {first_sentence or 'not stated'}."


def parse_file(path: Path, max_chars: int = 12000) -> Result:
    text, warnings = read_text(path)
    normalized = text.replace('\x00', ' ')
    doc_type = classify(normalized, path.name)
    if doc_type == 'invoice_or_quote':
        fields = extract_invoice_fields(normalized)
    elif doc_type == 'contract_or_legal_document':
        fields = extract_contract_fields(normalized)
    elif doc_type == 'resume_or_candidate_profile':
        fields = extract_resume_fields(normalized)
    elif doc_type == 'report_or_presentation':
        fields = extract_report_fields(normalized)
    else:
        fields = extract_generic_fields(normalized)
    if len(normalized.strip()) < 50:
        warnings.append('very little text was extracted; confirm whether the file is scanned, empty, or image-based')
    return Result(path.name, path.suffix.lower().lstrip('.'), doc_type, fields, summarize(doc_type, fields, normalized), warnings, limit(normalized, max_chars))


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument('file')
    ap.add_argument('--max-chars', type=int, default=12000)
    args = ap.parse_args()
    path = Path(args.file)
    if not path.exists():
        print(json.dumps({'error': f'file not found: {path}'}, indent=2))
        return 1
    try:
        print(json.dumps(asdict(parse_file(path, args.max_chars)), indent=2, ensure_ascii=False))
        return 0
    except Exception as exc:
        print(json.dumps({'error': str(exc), 'file': os.fspath(path)}, indent=2, ensure_ascii=False))
        return 1


if __name__ == '__main__':
    sys.exit(main())
